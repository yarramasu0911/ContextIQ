from docx import Document

try:
    import pdfplumber
except Exception:  # pragma: no cover
    pdfplumber = None

try:
    from PyPDF2 import PdfReader  # legacy fallback
except Exception:  # pragma: no cover
    PdfReader = None

# OCR for scanned PDFs is opt-in (requires Tesseract binary installed
# separately) — fall back silently if the deps or the binary aren't present.
try:
    import pytesseract
    from pdf2image import convert_from_path
except Exception:  # pragma: no cover
    pytesseract = None
    convert_from_path = None
import unicodedata
import re
import os
from bs4 import BeautifulSoup
import json
from striprtf.striprtf import rtf_to_text
import csv
from datetime import datetime
import numpy as np

# Optional imports — only needed for their respective formats. Don't fail at
# import time if the user hasn't installed them; fail when they try to upload
# that format.
try:
    from openpyxl import load_workbook  # .xlsx
except Exception:  # pragma: no cover
    load_workbook = None

try:
    from pptx import Presentation  # .pptx
except Exception:  # pragma: no cover
    Presentation = None

try:
    import yaml  # .yaml/.yml
except Exception:  # pragma: no cover
    yaml = None

try:
    import xml.etree.ElementTree as ET
except Exception:  # pragma: no cover
    ET = None


def clean_text(text: str) -> str:
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r"[\u200B-\u200D\uFEFF]", "", text)
    text = re.sub(r"-\n", "", text)
    text = re.sub(r"(?<!\n)\n(?!\n)", " ", text)
    text = re.sub(r"\n{2,}", "\n", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _ocr_pdf(file_path: str) -> str:
    """Best-effort OCR for scanned PDFs. Returns empty string if unavailable."""
    if pytesseract is None or convert_from_path is None:
        return ""
    try:
        images = convert_from_path(file_path, dpi=200)
    except Exception as e:
        print(f"[ocr] pdf2image failed ({e}); is poppler installed?")
        return ""
    parts = []
    for i, img in enumerate(images, start=1):
        try:
            text = pytesseract.image_to_string(img)
        except Exception as e:
            print(f"[ocr] tesseract failed on page {i}: {e}")
            continue
        if text.strip():
            parts.append(f"[Page {i}]\n{text}")
    return "\n\n".join(parts)


def extract_text_from_pdf(file_path: str) -> str:
    """Layout-aware PDF extraction via pdfplumber.

    - Extracts page text and tables (rendered as `cell | cell | cell` rows).
    - If the document has no extractable text (scanned PDF), falls back to OCR
      when pytesseract + poppler are available.
    """
    if pdfplumber is not None:
        parts = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, start=1):
                    page_parts = []
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        page_parts.append(page_text)

                    for table in page.extract_tables() or []:
                        rendered = []
                        for row in table:
                            cells = [
                                (c or "").replace("\n", " ").strip() for c in row
                            ]
                            if any(cells):
                                rendered.append(" | ".join(cells))
                        if rendered:
                            page_parts.append(
                                "[Table]\n" + "\n".join(rendered)
                            )

                    if page_parts:
                        parts.append(f"[Page {i}]\n" + "\n\n".join(page_parts))
        except Exception as e:
            print(f"[pdf] pdfplumber failed ({e}); falling back to PyPDF2")
            parts = []

        text = "\n\n".join(parts).strip()
        if text:
            return text

    # Fallback 1: PyPDF2
    if PdfReader is not None:
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            if text.strip():
                return text
        except Exception as e:
            print(f"[pdf] PyPDF2 failed: {e}")

    # Fallback 2: OCR (scanned PDF)
    ocr_text = _ocr_pdf(file_path)
    if ocr_text.strip():
        return ocr_text

    raise RuntimeError(
        f"Could not extract text from {file_path}. The PDF may be scanned "
        "(install pytesseract + Tesseract + poppler for OCR) or corrupted."
    )


def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_from_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])


def extract_from_html(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    return soup.get_text(separator="\n")


def extract_from_json(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    return json.dumps(data, indent=2)


def extract_from_rtf(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return rtf_to_text(f.read())


def extract_from_csv(file_path: str) -> str:
    rows = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))
    return "\n".join(rows)


def extract_from_tsv(file_path: str) -> str:
    rows = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        for line in f:
            rows.append(" | ".join(line.rstrip("\n").split("\t")))
    return "\n".join(rows)


def extract_from_xlsx(file_path: str) -> str:
    if load_workbook is None:
        raise RuntimeError(
            "openpyxl is required for .xlsx/.xls support. Install with "
            "`pip install openpyxl`."
        )
    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_from_pptx(file_path: str) -> str:
    if Presentation is None:
        raise RuntimeError(
            "python-pptx is required for .pptx support. Install with "
            "`pip install python-pptx`."
        )
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def extract_from_yaml(file_path: str) -> str:
    if yaml is None:
        # fall back to raw text — still searchable, just not pretty
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        data = yaml.safe_load(f)
    return json.dumps(data, indent=2, default=str)


def extract_from_xml(file_path: str) -> str:
    if ET is None:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
    except Exception:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    parts = []

    def walk(node, depth=0):
        indent = "  " * depth
        tag = node.tag.split("}")[-1]
        text = (node.text or "").strip()
        if text:
            parts.append(f"{indent}{tag}: {text}")
        else:
            parts.append(f"{indent}{tag}")
        for child in node:
            walk(child, depth + 1)

    walk(root)
    return "\n".join(parts)


def split_into_sentences(text: str) -> list:
    """Split text into sentences."""
    sentences = re.split(r'(?<=[.!?;:])\s+', text)
    return [s.strip() for s in sentences if len(s.strip()) > 10]


def semantic_chunking(text: str, embedding_model=None,
                      max_chunk_size: int = 500,
                      similarity_threshold: float = 0.45) -> list:
    """Split text into chunks based on meaning similarity."""

    sentences = split_into_sentences(text)

    if len(sentences) <= 1:
        return [text.strip()] if text.strip() else []

    if embedding_model is None:
        from sentence_transformers import SentenceTransformer
        embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

    embeddings = embedding_model.encode(
        sentences, convert_to_numpy=True, normalize_embeddings=True
    )

    chunks = []
    current_chunk = [sentences[0]]
    current_length = len(sentences[0])

    for i in range(1, len(sentences)):
        similarity = float(np.dot(embeddings[i - 1], embeddings[i]))

        should_split = False

        if similarity < similarity_threshold:
            should_split = True

        if current_length + len(sentences[i]) > max_chunk_size:
            should_split = True

        if should_split and current_chunk:
            chunks.append(" ".join(current_chunk))
            current_chunk = [sentences[i]]
            current_length = len(sentences[i])
        else:
            current_chunk.append(sentences[i])
            current_length += len(sentences[i])

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    # Merge very small chunks with previous
    merged = []
    for chunk in chunks:
        if merged and len(chunk) < 100:
            merged[-1] = merged[-1] + " " + chunk
        else:
            merged.append(chunk)

    return merged


def split_into_chunks(text: str, chunk_size: int = 500, overlap: int = 100) -> list:
    """Fallback: fixed-size chunking with overlap."""
    if not text:
        return []
    if overlap >= chunk_size:
        raise ValueError("overlap must be smaller than chunk_size")

    chunks = []
    start = 0
    step = chunk_size - overlap
    text_length = len(text)

    while start < text_length:
        end = min(start + chunk_size, text_length)
        if end < text_length:
            space_pos = text.rfind(" ", start, end)
            if space_pos > start:
                end = space_pos
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += step

    return chunks


def ingest_document(file_path: str) -> dict:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext in [".txt", ".md", ".log"]:
        text = extract_text_from_txt(file_path)
    elif ext in [".docx", ".doc"]:
        text = extract_from_docx(file_path)
    elif ext == ".csv":
        text = extract_from_csv(file_path)
    elif ext == ".tsv":
        text = extract_from_tsv(file_path)
    elif ext in [".html", ".htm"]:
        text = extract_from_html(file_path)
    elif ext == ".json":
        text = extract_from_json(file_path)
    elif ext == ".rtf":
        text = extract_from_rtf(file_path)
    elif ext in [".xlsx", ".xls"]:
        text = extract_from_xlsx(file_path)
    elif ext in [".pptx", ".ppt"]:
        text = extract_from_pptx(file_path)
    elif ext in [".yaml", ".yml"]:
        text = extract_from_yaml(file_path)
    elif ext == ".xml":
        text = extract_from_xml(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    text = clean_text(text)
    return {
        "content": text,
        "metadata": {
            "file_name": os.path.basename(file_path),
            "file_type": ext,
            "source": "upload",
            "ingested_at": datetime.now().isoformat(),
        },
    }


# Per-extension chunk size tuning. Tables/spreadsheets/slides need bigger
# windows; prose is fine at 500 chars.
_CHUNK_SIZE_BY_EXT = {
    ".csv": 1000, ".tsv": 1000,
    ".xlsx": 1000, ".xls": 1000,
    ".pptx": 800, ".ppt": 800,
    ".json": 1200, ".yaml": 1200, ".yml": 1200, ".xml": 1200,
    ".log": 800,
    ".html": 600, ".htm": 600,
}


def process_document(
    file_path: str,
    chunk_size: int | None = None,
    overlap: int = 100,
    use_semantic: bool = True,
    embedding_model=None,
    parent_window: int = 1,
) -> list:
    """Produce chunks for indexing, each carrying a `parent_content` field
    built from [chunk_{i-parent_window}, ..., chunk_{i+parent_window}] — the
    large context we show the LLM when this chunk is retrieved.
    """
    document = ingest_document(file_path)
    ext = document["metadata"]["file_type"]

    if chunk_size is None:
        chunk_size = _CHUNK_SIZE_BY_EXT.get(ext, 500)

    if use_semantic:
        chunks = semantic_chunking(
            document["content"],
            embedding_model=embedding_model,
            max_chunk_size=chunk_size,
            similarity_threshold=0.45,
        )
    else:
        chunks = split_into_chunks(
            document["content"], chunk_size=chunk_size, overlap=overlap
        )

    chunk_docs = []
    total_chunks = len(chunks)

    for i, chunk in enumerate(chunks):
        lo = max(0, i - parent_window)
        hi = min(total_chunks, i + parent_window + 1)
        parent_content = " ".join(chunks[lo:hi])
        chunk_docs.append(
            {
                "content": chunk,
                "parent_content": parent_content,
                "metadata": {
                    **document["metadata"],
                    "chunk_index": i,
                    "total_chunks": total_chunks,
                    "chunking_method": "semantic" if use_semantic else "fixed",
                    "chunk_size_used": chunk_size,
                    "parent_window": parent_window,
                },
            }
        )

    return chunk_docs